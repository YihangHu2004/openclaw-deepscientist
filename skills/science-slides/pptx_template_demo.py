# -*- coding: utf-8 -*-
"""
Demo: Build a PPT from a user template (.pptx).

Usage:
  1. Set TEMPLATE_PATH to your .pptx template (or ~/slides/templates/*.pptx).
  2. Set OUTPUT_PATH for the generated file.
  3. Run: python pptx_template_demo.py

Key techniques:
  - clone_slide()      : copies a prototype slide preserving background images
                         (rId relationship remapping)
  - set_placeholder()  : writes text into a placeholder, clearing ALL paragraphs
                         first to avoid template residue text
  - template caching   : copies template to project slides/ on first run
"""
import os, sys, shutil
sys.stdout.reconfigure(encoding='utf-8')
from lxml import etree
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

# ── Paths (edit these) ──────────────────────────────────────────────────────
TEMPLATE_PATH = r'~/slides/templates/your_template.pptx'   # user template
OUTPUT_PATH   = r'state/projects/<slug>/slides/output.pptx'
CACHE_DIR     = r'state/projects/<slug>/slides'
CACHE_PATH    = os.path.join(CACHE_DIR, 'template_cache.pptx')

# ── Theme colors (adjust to match your template) ────────────────────────────
ORANGE    = RGBColor(0xED, 0x7D, 0x31)
DARK_BLUE = RGBColor(0x44, 0x54, 0x6A)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)


# ── Core helpers ─────────────────────────────────────────────────────────────

def clone_slide(prs, source_slide):
    """Clone a slide, preserving background images via rId remapping."""
    new_slide = prs.slides.add_slide(source_slide.slide_layout)

    rId_map = {}
    for rId, rel in source_slide.part.rels.items():
        if 'slideLayout' in rel.reltype:
            continue
        if not rel.is_external:
            new_rId = new_slide.part.relate_to(rel.target_part, rel.reltype)
            rId_map[rId] = new_rId

    src_xml = etree.tostring(source_slide.shapes._spTree, encoding='unicode')
    for old_rId, new_rId in rId_map.items():
        src_xml = src_xml.replace(f'"{old_rId}"', f'"{new_rId}"')

    new_spTree = etree.fromstring(src_xml)
    dst_spTree = new_slide.shapes._spTree
    dst_spTree.getparent().replace(dst_spTree, new_spTree)
    return new_slide


def set_placeholder(slide, idx, text, size_pt=None, bold=None, color=None):
    """Write text into placeholder idx, removing all template residue paragraphs."""
    from pptx.oxml.ns import qn
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            tf = ph.text_frame
            tf.word_wrap = True
            txBody = tf._txBody
            for extra_p in txBody.findall(qn('a:p'))[1:]:
                txBody.remove(extra_p)
            p = tf.paragraphs[0]
            p.clear()
            run = p.add_run()
            run.text = text
            if size_pt:
                run.font.size = Pt(size_pt)
            if bold is not None:
                run.font.bold = bold
            if color:
                run.font.color.rgb = color
            return


def set_bullets(slide, body_idx, items, size_pt=20):
    """Write a bullet list into a body placeholder."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == body_idx:
            tf = ph.text_frame
            tf.word_wrap = True
            for i, text in enumerate(items):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.clear()
                run = p.add_run()
                run.text = text
                run.font.size = Pt(size_pt)
                run.font.name = "Microsoft YaHei"
                p.level = 0
            return


# ── Step 0: Template caching ─────────────────────────────────────────────────
os.makedirs(CACHE_DIR, exist_ok=True)
if not os.path.exists(CACHE_PATH):
    shutil.copy2(TEMPLATE_PATH, CACHE_PATH)
    print(f"模板已缓存：{CACHE_PATH}")
else:
    print(f"使用已缓存模板：{CACHE_PATH}")

prs = Presentation(CACHE_PATH)
print(f"Template loaded: {len(prs.slides)} demo slides, {len(prs.slide_layouts)} layouts")

# Prototype slides (slide[0] = cover, slide[1] = content with background)
COVER_PROTO   = prs.slides[0]
CONTENT_PROTO = prs.slides[1]
N_PROTOS      = len(prs.slides)

# ── Slides ───────────────────────────────────────────────────────────────────

cover = clone_slide(prs, COVER_PROTO)
set_placeholder(cover, 0, "研究题目\n副标题", size_pt=36, bold=True)
set_placeholder(cover, 1, "汇报人：XXX  |  导师：XXX  |  日期")

s2 = clone_slide(prs, CONTENT_PROTO)
set_placeholder(s2, 0, "研究背景", size_pt=28, bold=True)
set_bullets(s2, 1, ["背景点 1", "背景点 2", "背景点 3"])

# Add more slides by calling clone_slide(prs, CONTENT_PROTO) + set_placeholder...

# ── Remove prototype slides ───────────────────────────────────────────────────
for _ in range(N_PROTOS):
    rId = prs.slides._sldIdLst[0].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[0]

print(f"Removed {N_PROTOS} demo slides. Final: {len(prs.slides)} slides")

prs.save(OUTPUT_PATH)
print(f"Saved → {OUTPUT_PATH}")
